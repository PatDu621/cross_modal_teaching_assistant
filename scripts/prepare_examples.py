#!/usr/bin/env python
# scripts/prepare_examples.py
"""
自动生成测试样例文件（sample.pdf 和 sample.pptx）
放到 examples/ 目录下，供 processor 单元测试使用
"""

import os
import sys

# 添加项目根目录到 Python 路径
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

from modules.utils.helpers import ensure_dir


def create_dummy_image(image_path: str, size=(300, 200), color=(100, 150, 200)):
    """用 PIL 生成一张纯色测试图片"""
    try:
        from PIL import Image
        img = Image.new('RGB', size, color=color)
        img.save(image_path)
        print(f"  ✓ 生成图片: {os.path.basename(image_path)}")
        return True
    except ImportError:
        print("  ✗ 未安装 Pillow，请运行: pip install Pillow")
        return False


def create_sample_pptx(pptx_path: str):
    """生成包含 4 页（文本+图片）的测试 PPT"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  ✗ 未安装 python-pptx，请运行: pip install python-pptx")
        return False

    prs = Presentation()
    
    # ---- 第1页：标题 ----
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "跨模态教学辅助系统 - 测试材料"
    subtitle.text = "用于验证多模态输入处理模块\n\n生成日期: 2026年7月"

    # ---- 第2页：文本列表 ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "牛顿第二定律核心要点"
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.text = "• 力是改变物体运动状态的原因"
    p = text_frame.add_paragraph()
    p.text = "• 加速度与合外力成正比 (F = ma)"
    p.level = 0
    p = text_frame.add_paragraph()
    p.text = "• 加速度与质量成反比"
    p.level = 0
    p = text_frame.add_paragraph()
    p.text = "• 适用于惯性参考系"
    p.level = 0

    # ---- 第3页：纯图片 ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 生成临时图片并插入
    temp_img = os.path.join(os.path.dirname(pptx_path), "_temp_force_diagram.png")
    if create_dummy_image(temp_img, size=(400, 300), color=(200, 100, 100)):
        left = Inches(1.5)
        top = Inches(1.5)
        slide.shapes.add_picture(temp_img, left, top, height=Inches(4.5))
        os.remove(temp_img)
        # 添加文字说明
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(1))
        tx_box.text = "图1：力、质量、加速度关系示意图"
        tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---- 第4页：文字 + 图片混合 ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 左侧文字
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(4))
    text_frame = tx_box.text_frame
    text_frame.text = "【混合页面】\n\n光反应发生在类囊体膜上，将光能转化为 ATP 和 NADPH。暗反应在叶绿体基质中进行，利用 ATP 和 NADPH 将 CO₂ 固定为糖类。"
    text_frame.paragraphs[0].font.size = 240000  # 24pt
    # 右侧图片
    temp_img2 = os.path.join(os.path.dirname(pptx_path), "_temp_photosynthesis.png")
    if create_dummy_image(temp_img2, size=(300, 250), color=(100, 200, 150)):
        slide.shapes.add_picture(temp_img2, Inches(5), Inches(1.5), height=Inches(3.5))
        os.remove(temp_img2)

    # 保存
    prs.save(pptx_path)
    print(f"  ✓ 生成 PPT: {os.path.basename(pptx_path)} (4页)")
    return True


def create_sample_pdf(pdf_path: str):
    """使用 PyMuPDF 生成包含文本和内嵌图片的测试 PDF"""
    try:
        import fitz  # PyMuPDF
        from PIL import Image
        import io
    except ImportError as e:
        print(f"  ✗ 缺少依赖: {e}")
        print("    请运行: pip install PyMuPDF Pillow")
        return False

    doc = fitz.open()

    # ---- 第1页：纯文本 ----
    page = doc.new_page(width=595, height=842)  # A4
    page.insert_text((72, 100), "跨模态教学辅助系统 - 测试文档", fontsize=20, color=(0, 0, 1))
    page.insert_text((72, 150), "=" * 40, fontsize=14)
    page.insert_text((72, 190), "本PDF用于测试多模态输入处理模块的文本和图片提取功能。", fontsize=12)
    page.insert_text((72, 220), "包含以下内容：", fontsize=12)
    page.insert_text((72, 250), "1. 文字段落（本页）", fontsize=12)
    page.insert_text((72, 275), "2. 项目符号列表（下一页）", fontsize=12)
    page.insert_text((72, 300), "3. 内嵌图片（第三页）", fontsize=12)

    # ---- 第2页：列表 ----
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 100), "量子力学基础概念", fontsize=18, color=(0.5, 0, 0.5))
    items = [
        "• 波粒二象性：微观粒子兼具波动性和粒子性",
        "• 不确定性原理：Δx · Δp ≥ ħ/2",
        "• 薛定谔方程：描述量子态随时间演化",
        "• 量子纠缠：多个粒子之间存在非经典关联",
    ]
    y = 150
    for item in items:
        page.insert_text((72, y), item, fontsize=12)
        y += 30

    # ---- 第3页：内嵌图片 ----
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 80), "图：原子结构示意图", fontsize=14, color=(0, 0.5, 0))

    # 生成一张彩色圆点图作为测试图片
    img_pil = Image.new('RGB', (400, 300), color=(240, 240, 240))
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img_pil)
    # 画三个圆表示原子核+电子
    draw.ellipse([150, 100, 250, 200], fill=(255, 0, 0), outline=(0, 0, 0))
    draw.ellipse([50, 50, 100, 100], fill=(0, 0, 255), outline=(0, 0, 0))
    draw.ellipse([300, 50, 350, 100], fill=(0, 0, 255), outline=(0, 0, 0))
    draw.ellipse([100, 200, 150, 250], fill=(0, 0, 255), outline=(0, 0, 0))
    draw.ellipse([280, 200, 330, 250], fill=(0, 0, 255), outline=(0, 0, 0))
    # 标注
    draw.text((180, 220), "原子核", fill=(0, 0, 0))
    draw.text((30, 80), "电子", fill=(0, 0, 0))

    # 将 PIL 图片转为字节插入 PDF
    img_bytes = io.BytesIO()
    img_pil.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    page.insert_image((72, 110, 72+400, 110+300), stream=img_bytes)

    # 保存 PDF
    doc.save(pdf_path)
    doc.close()
    print(f"  ✓ 生成 PDF: {os.path.basename(pdf_path)} (3页)")
    return True


def main():
    print("=" * 50)
    print("准备测试样例文件")
    print("=" * 50)

    examples_dir = os.path.join(PROJECT_ROOT, "examples")
    ensure_dir(examples_dir)

    # 1. 生成 PPT
    pptx_path = os.path.join(examples_dir, "sample.pptx")
    if os.path.exists(pptx_path):
        print(f"  ⚠️ 文件已存在，跳过: {os.path.basename(pptx_path)}")
    else:
        create_sample_pptx(pptx_path)

    # 2. 生成 PDF
    pdf_path = os.path.join(examples_dir, "sample.pdf")
    if os.path.exists(pdf_path):
        print(f"  ⚠️ 文件已存在，跳过: {os.path.basename(pdf_path)}")
    else:
        create_sample_pdf(pdf_path)

    # 3. 生成一张额外图片（方便单独测试图片解析）
    img_path = os.path.join(examples_dir, "sample.png")
    if not os.path.exists(img_path):
        create_dummy_image(img_path, size=(200, 200), color=(150, 200, 100))
        print(f"  ✓ 生成图片: sample.png")

    print("\n" + "=" * 50)
    print("✅ 样例文件准备完成！")
    print(f"目录: {examples_dir}")
    print("现在可以运行单元测试: python -m pytest tests/ -v")
    print("或直接测试单个文件: python -m modules.input_processor.processor examples/sample.pdf")
    print("=" * 50)


if __name__ == "__main__":
    main()