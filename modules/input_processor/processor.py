"""
processor.py - 多模态输入处理模块
负责解析用户上传的教学材料（PPT/PDF/图片/视频），提取文本和图像路径，
输出统一的结构化JSON数据供后续模型调用使用。

支持的输入格式：
  - PDF (.pdf)   → 提取每页文本 + 内嵌图片（使用 PyMuPDF）
  - PPT (.pptx)  → 提取每页文本 + 内嵌图片（使用 python-pptx）
  - 图片 (.png/.jpg等) → 直接作为图像路径保留
  - 视频 (.mp4/.avi等) → 抽帧保存为图片（可选）
"""

import os
import shutil
from typing import List, Dict, Any, Tuple

# PPT解析
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 图像处理
from PIL import Image

# 项目内工具
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))
from modules.utils.helpers import (
    generate_material_id,
    get_file_type,
    get_file_name_without_ext,
    ensure_dir,
    save_json,
)


class MaterialProcessor:
    """
    教学材料处理器
    负责将各种格式的输入文件统一解析为结构化的材料数据
    """

    def __init__(self, raw_dir: str = "data/raw",
                 processed_dir: str = "data/processed"):
        """
        初始化处理器
        Args:
            raw_dir: 原始上传文件的存放目录
            processed_dir: 处理结果的输出目录
        """
        self.raw_dir = os.path.abspath(raw_dir)
        self.processed_dir = os.path.abspath(processed_dir)
        ensure_dir(self.raw_dir)
        ensure_dir(self.processed_dir)

    def process(self, file_path: str) -> Dict[str, Any]:
        """
        处理上传的文件，自动判断类型并调用对应的解析方法
        Args:
            file_path: 用户上传的文件路径
        Returns:
            统一格式的材料数据字典
        """
        file_path = os.path.abspath(file_path)
        file_type = get_file_type(file_path)
        material_id = generate_material_id()
        title = get_file_name_without_ext(file_path)

        # 将原始文件复制到 raw 目录
        raw_copy_path = os.path.join(self.raw_dir, os.path.basename(file_path))
        ensure_dir(self.raw_dir)
        if not os.path.exists(raw_copy_path):
            shutil.copy2(file_path, raw_copy_path)

        # 根据文件类型选择解析策略
        if file_type == "pdf":
            text_blocks, image_paths = self._parse_pdf(file_path, material_id)
        elif file_type == "ppt":
            text_blocks, image_paths = self._parse_ppt(file_path, material_id)
        elif file_type == "image":
            text_blocks, image_paths = self._parse_image(file_path, material_id)
        elif file_type == "video":
            text_blocks, image_paths = self._parse_video(file_path, material_id)
        else:
            # 不支持的文件格式，返回空结果
            print(f"[警告] 不支持的文件格式: {file_path}")
            text_blocks, image_paths = [], []

        # 组装统一的输出结构
        material_data = {
            "material_id": material_id,
            "title": title,
            "text_blocks": text_blocks,
            "image_paths": image_paths,
        }

        # 保存处理结果到 processed 目录
        output_path = os.path.join(self.processed_dir, f"{material_id}.json")
        save_json(material_data, output_path)
        print(f"[完成] 材料已处理并保存到: {output_path}")

        return material_data

    # ============================================================
    # PDF 解析（使用 PyMuPDF）
    # ============================================================
    def _parse_pdf(self, file_path: str, material_id: str
                   ) -> Tuple[List[Dict], List[str]]:
        """
        解析 PDF 文件，提取每页文本和内嵌图片
        Args:
            file_path: PDF 文件路径
            material_id: 材料ID
        Returns:
            (text_blocks, image_paths)
        """
        text_blocks = []
        image_paths = []

        # 创建该材料专属的图片目录
        img_dir = os.path.join(self.processed_dir, f"{material_id}_images")
        ensure_dir(img_dir)

        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc, start=1):
                # ---- 提取文本 ----
                text = page.get_text().strip()
                if text:
                    text_blocks.append({"page": page_num, "text": text})

                # ---- 提取图片 ----
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]  # e.g., 'png', 'jpeg'

                        # 生成图片文件名
                        img_filename = f"page{page_num}_img{img_index}.{image_ext}"
                        img_path = os.path.join(img_dir, img_filename)

                        # 保存图片
                        with open(img_path, "wb") as f:
                            f.write(image_bytes)
                        image_paths.append(os.path.abspath(img_path))
                        print(f"[提取图片] {img_filename} (PDF页码{page_num})")
                    except Exception as e:
                        # 单张图片提取失败不影响整体
                        print(f"[跳过图片] PDF页码{page_num}中的图片提取失败: {e}")

            doc.close()
            if not text_blocks:
                print("[提示] PDF 中未提取到文本内容，可能是扫描版 PDF")

        except ImportError:
            print("[错误] 缺少 PyMuPDF 依赖，请运行: pip install PyMuPDF")
        except Exception as e:
            print(f"[错误] PDF 解析失败: {e}")

        return text_blocks, image_paths

    # ============================================================
    # PPT 解析（增强版）
    # ============================================================
    def _parse_ppt(self, file_path: str, material_id: str
                   ) -> Tuple[List[Dict], List[str]]:
        """
        解析 PPT 文件，提取每页文本和内嵌图片
        Args:
            file_path: PPT 文件路径
            material_id: 材料ID
        Returns:
            (text_blocks, image_paths)
        """
        text_blocks = []
        image_paths = []

        img_dir = os.path.join(self.processed_dir, f"{material_id}_images")
        ensure_dir(img_dir)

        try:
            prs = Presentation(file_path)

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_texts = []

                for shape in slide.shapes:
                    # ---- 提取文本 ----
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_texts.append(para_text)

                    # ---- 提取内嵌图片 ----
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            ext = image.content_type.split("/")[-1]
                            if ext == "jpeg":
                                ext = "jpg"
                            img_filename = f"slide{slide_idx}_shape{shape.shape_id}.{ext}"
                            img_path = os.path.join(img_dir, img_filename)

                            with open(img_path, "wb") as f:
                                f.write(image.blob)
                            image_paths.append(os.path.abspath(img_path))
                            print(f"[提取图片] {img_filename} (幻灯片{slide_idx})")
                        except Exception as e:
                            print(f"[跳过图片] 幻灯片{slide_idx}中的图片提取失败: {e}")

                if slide_texts:
                    text_blocks.append({
                        "page": slide_idx,
                        "text": "\n".join(slide_texts),
                    })

            if not text_blocks:
                print("[提示] PPT 中未提取到文本内容（可能全是图片的PPT）")

        except Exception as e:
            print(f"[错误] PPT 解析失败: {e}")

        # 注：PPT整页渲染为图片需要LibreOffice等外部工具，暂不实现。
        return text_blocks, image_paths

    # ============================================================
    # 图片解析
    # ============================================================
    def _parse_image(self, file_path: str, material_id: str
                     ) -> Tuple[List[Dict], List[str]]:
        """
        处理用户上传的图片文件
        """
        img_dir = os.path.join(self.processed_dir, f"{material_id}_images")
        ensure_dir(img_dir)

        ext = os.path.splitext(file_path)[1]
        dest_path = os.path.join(img_dir, f"uploaded_image{ext}")
        shutil.copy2(file_path, dest_path)
        abs_dest = os.path.abspath(dest_path)

        # 验证图片是否有效
        try:
            with Image.open(abs_dest) as img:
                print(f"[图片信息] 尺寸={img.size}, 格式={img.format}")
            if img.width > 4000 or img.height > 4000:
                print("[提示] 图片尺寸较大，建议压缩后上传")
        except Exception as e:
            print(f"[错误] 图片文件无效: {e}")
            return [], []

        return [], [abs_dest]

    # ============================================================
    # 视频抽帧（新增）
    # ============================================================
    def _parse_video(self, file_path: str, material_id: str,
                     frame_interval: float = 1.0) -> Tuple[List[Dict], List[str]]:
        """
        从视频中抽取关键帧作为图片输入
        Args:
            file_path: 视频文件路径
            material_id: 材料ID
            frame_interval: 抽帧间隔（秒），默认每秒1帧
        Returns:
            (text_blocks, image_paths)  text_blocks为空
        """
        image_paths = []
        img_dir = os.path.join(self.processed_dir, f"{material_id}_images")
        ensure_dir(img_dir)

        try:
            import cv2

            cap = cv2.VideoCapture(file_path)
            if not cap.isOpened():
                print(f"[错误] 无法打开视频文件: {file_path}")
                return [], []

            fps = cap.get(cv2.CAP_PROP_FPS)
            if fps <= 0:
                fps = 30  # 后备
            frame_skip = int(fps * frame_interval)
            if frame_skip < 1:
                frame_skip = 1

            frame_count = 0
            saved_count = 0
            while True:
                ret, frame = cap.read()
                if not ret:
                    break
                if frame_count % frame_skip == 0:
                    img_filename = f"frame_{saved_count:04d}.jpg"
                    img_path = os.path.join(img_dir, img_filename)
                    cv2.imwrite(img_path, frame)
                    image_paths.append(os.path.abspath(img_path))
                    saved_count += 1
                    print(f"[提取帧] {img_filename}")
                frame_count += 1

            cap.release()
            print(f"[完成] 从视频中抽取 {saved_count} 帧")

        except ImportError:
            print("[错误] 缺少 opencv-python 依赖，请运行: pip install opencv-python")
        except Exception as e:
            print(f"[错误] 视频处理失败: {e}")

        return [], image_paths


# ============================================================
# 简单测试入口
# ============================================================
if __name__ == "__main__":
    import sys

    # 创建处理器实例，使用项目默认路径
    processor = MaterialProcessor(
        raw_dir="../../data/raw",
        processed_dir="../../data/processed",
    )

    if len(sys.argv) > 1:
        test_file = sys.argv[1]
    else:
        print("用法: python processor.py <文件路径>")
        print("示例: python processor.py ../examples/sample.pdf")
        sys.exit(1)

    result = processor.process(test_file)
    print("\n===== 处理结果 =====")
    print(f"材料ID: {result['material_id']}")
    print(f"标题: {result['title']}")
    print(f"文本块数量: {len(result['text_blocks'])}")
    print(f"图片数量: {len(result['image_paths'])}")
    for block in result["text_blocks"]:
        print(f"  - 第{block['page']}页: {block['text'][:80]}...")
    for img in result["image_paths"]:
        print(f"  - 图片: {img}")
